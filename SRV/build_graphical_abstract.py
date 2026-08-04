#!/usr/bin/env python3
"""Build the editable graphical abstract for the revised manuscript.

The artwork is composed only of PowerPoint vector shapes and text. It does not
use generative AI imagery or third-party artwork.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as Shape
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


W, H = Inches(13.333), Inches(5.333)
NAVY, BLUE, TEAL, ORANGE, RED, GREY, PALE_BLUE, PALE_GREEN, PALE_RED = (
    "17324D", "3978B8", "0B745F", "DF7B22", "B95454", "59636E", "E8F1FA", "E7F6EF", "FBEDEE"
)


def rgb(hexcode: str) -> RGBColor:
    return RGBColor.from_string(hexcode)


def add_text(slide, x, y, w, h, text, size, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text; run.font.name = "Arial"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = rgb(color)
    return box


def card(slide, x, y, w, h, fill, title, body, accent):
    shape = slide.shapes.add_shape(Shape.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(accent); shape.line.width = Pt(1.4)
    strip = slide.shapes.add_shape(Shape.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(h))
    strip.fill.solid(); strip.fill.fore_color.rgb = rgb(accent); strip.line.fill.background()
    add_text(slide, x + 0.22, y + 0.14, w - 0.38, 0.34, title, 14, accent, True)
    add_text(slide, x + 0.22, y + 0.58, w - 0.38, h - 0.72, body, 11, NAVY, False)


def arrow(slide, x, y, w, label):
    shape = slide.shapes.add_shape(Shape.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.38))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(ORANGE); shape.line.fill.background()
    add_text(slide, x - 0.30, y - 0.52, w + 0.60, 0.44, label, 9.5, "8A4D17", True, PP_ALIGN.CENTER)


def build(output: Path):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(255, 255, 255)

    add_text(slide, 0.42, 0.18, 12.45, 0.45, "Proposed tBLG coatings for hard-carbon sodium-ion anodes", 22, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, 0.42, 0.61, 12.45, 0.24, "Evidence-gated assessment; no validated battery transfer function", 11, GREY, False, PP_ALIGN.CENTER)

    card(slide, 0.45, 1.05, 3.45, 2.95, PALE_BLUE,
         "PUBLISHED LEVEL-0 EVIDENCE",
         "Planar twisted bilayer graphene\nAqueous outer-sphere redox probes\nStrong, probe-dependent twist response", BLUE)
    # Simple layered-graphene icon, built from editable shapes.
    for yy, col, rotation in [(3.18, BLUE, -7), (3.38, TEAL, 7)]:
        layer = slide.shapes.add_shape(Shape.PARALLELOGRAM, Inches(1.00), Inches(yy), Inches(1.70), Inches(0.32))
        layer.fill.solid(); layer.fill.fore_color.rgb = rgb(col); layer.fill.transparency = 20
        layer.line.color.rgb = rgb(NAVY); layer.line.width = Pt(0.9)
        layer.rotation = rotation
    add_text(slide, 0.96, 3.70, 1.84, 0.20, "relative twist θ", 10, NAVY, True, PP_ALIGN.CENTER)

    arrow(slide, 4.03, 2.28, 1.02, "no validated\ntransfer function")

    card(slide, 5.17, 1.05, 3.08, 2.95, PALE_RED,
         "EVIDENCE GAP",
         "No relevant non-aqueous replication identified\n\nNo tBLG-coated hard-carbon sodium-ion cell test or powder-coating route\n\nNo process, commercial, or policy evidence", RED)

    arrow(slide, 8.38, 2.28, 1.08, "advance only after\ngate criteria are met")

    card(slide, 9.58, 1.05, 3.29, 2.95, PALE_GREEN,
         "NEXT DECISION: EVIDENCE GATES",
         "1  Relevant-electrolyte model with controls\n2  Powder coverage, twist and adhesion\n3  Endpoint-specific half-cell effects\n\nThen: full cell → process inventory/LCA → stakeholder evidence", TEAL)

    footer = slide.shapes.add_shape(Shape.ROUNDED_RECTANGLE, Inches(0.80), Inches(4.33), Inches(11.73), Inches(0.54))
    footer.fill.solid(); footer.fill.fore_color.rgb = rgb("FFF5E8"); footer.line.color.rgb = rgb(ORANGE); footer.line.width = Pt(0.8)
    add_text(slide, 0.95, 4.41, 11.42, 0.27, "Level-0 evidence alone supports no conclusions about battery performance, cost competitiveness, commercialization, or policy.", 10.5, "8A4D17", True, PP_ALIGN.CENTER)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(); parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args(); build(args.output)
