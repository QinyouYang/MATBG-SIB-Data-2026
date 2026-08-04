#!/usr/bin/env python3
"""Validate the revised Journal of Energy Storage submission package.

This validator checks document structure, citation/reference bijection and first
appearance, recent-reference coverage, three-line tables, figure binaries,
headline limits, data consistency, disclosure sections, and cross-file titles.
It does not replace visual page inspection or publication of the pending Zenodo
version.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from zipfile import ZipFile

import pandas as pd
from PIL import Image
from docx import Document
from docx.document import Document as _Document
from docx.oxml.ns import qn
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PACKAGE = ROOT / "Zenodo_Package"
MANUSCRIPT = ROOT / "Manuscript.docx"
HIGHLIGHTS = ROOT / "Highlights.docx"
COVER = ROOT / "Cover Letter.docx"
GRAPHICAL_ABSTRACT = ROOT / "图文摘要.pptx"
RESPONSE = ROOT / "Response_to_Reviewers.docx"
REPORT = ROOT / "Submission_Validation_Report.json"
EXPECTED_TITLE = (
    "An Evidence-Gated Assessment of Proposed Twisted Bilayer Graphene "
    "Coatings for Hard-Carbon Sodium-Ion Anodes"
)


def iter_block_items(doc: _Document):
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def citation_numbers(text: str):
    for payload in re.findall(r"\[([0-9\s,;\-–]+)\]", text):
        for token in re.split(r"[,;]", payload):
            token = token.strip().replace("–", "-")
            if not re.fullmatch(r"\d+(?:-\d+)?", token):
                continue
            if "-" in token:
                low, high = map(int, token.split("-"))
                yield from range(low, high + 1)
            else:
                yield int(token)


def all_text(doc: _Document) -> str:
    values = [p.text for p in doc.paragraphs]
    values.extend(cell.text for table in doc.tables for row in table.rows for cell in row.cells)
    return "\n".join(values)


def pptx_text(path: Path) -> str:
    prs = Presentation(path)
    values = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                values.append(shape.text)
    return "\n".join(values)


def check_three_line(table: Table, name: str, errors: list[str]) -> None:
    title_row = table.rows[0].cells[0].text.strip().lower() == "nomenclature"
    header_index = 1 if title_row else 0
    last_index = len(table.rows) - 1
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            tc_pr = cell._tc.get_or_add_tcPr()
            if tc_pr.find(qn("w:shd")) is not None:
                errors.append(f"{name}: cell shading remains")
            borders = tc_pr.find(qn("w:tcBorders"))
            if borders is None:
                errors.append(f"{name}: missing explicit borders")
                continue
            for edge in ("left", "right", "insideH", "insideV"):
                node = borders.find(qn(f"w:{edge}"))
                if node is not None and node.get(qn("w:val")) not in ("nil", "none"):
                    errors.append(f"{name}: non-three-line border {edge} remains")
            top = borders.find(qn("w:top"))
            bottom = borders.find(qn("w:bottom"))
            top_on = top is not None and top.get(qn("w:val")) == "single"
            bottom_on = bottom is not None and bottom.get(qn("w:val")) == "single"
            if row_index == 0 and not top_on:
                errors.append(f"{name}: top rule missing")
            if row_index == header_index and not bottom_on:
                errors.append(f"{name}: header-bottom rule missing")
            if row_index == last_index and not bottom_on:
                errors.append(f"{name}: final-bottom rule missing")
            if row_index not in (header_index, last_index) and bottom_on:
                errors.append(f"{name}: extra interior bottom rule at row {row_index}")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--report", type=Path, default=REPORT)
    args = parser.parse_args()
    errors: list[str] = []
    warnings: list[str] = []
    metrics: dict[str, object] = {}

    required = [MANUSCRIPT, HIGHLIGHTS, COVER, GRAPHICAL_ABSTRACT, RESPONSE]
    for path in required:
        if not path.exists():
            errors.append(f"Missing required file: {path.name}")
    if errors:
        args.report.write_text(json.dumps({"errors": errors}, indent=2), encoding="utf-8")
        return 1

    doc = Document(MANUSCRIPT)
    text = all_text(doc)
    if doc.paragraphs[0].text != EXPECTED_TITLE:
        errors.append("Manuscript title does not match the approved evidence-gated title")

    abstract_heading = next(i for i, p in enumerate(doc.paragraphs) if p.text == "Abstract")
    abstract = doc.paragraphs[abstract_heading + 1].text
    abstract_words = re.findall(r"\b[\w+–-]+\b", abstract)
    metrics["abstract_words"] = len(abstract_words)
    if len(abstract_words) > 250:
        errors.append(f"Abstract exceeds 250 words: {len(abstract_words)}")

    keywords = next(p.text for p in doc.paragraphs if p.text.startswith("Keywords:"))
    keyword_count = len([part for part in keywords.split(":", 1)[1].split(";") if part.strip()])
    metrics["keyword_count"] = keyword_count
    if not 1 <= keyword_count <= 7:
        errors.append(f"Keyword count outside 1–7: {keyword_count}")

    ref_heading = next(i for i, p in enumerate(doc.paragraphs) if p.text == "References")
    appendix_heading = next(i for i, p in enumerate(doc.paragraphs) if p.text == "Appendix")
    refs = doc.paragraphs[ref_heading + 1 : appendix_heading]
    metrics["reference_count"] = len(refs)
    if len(refs) < 30:
        errors.append(f"Fewer than 30 focused references: {len(refs)}")

    ordered: list[int] = []
    cited: set[int] = set()
    for block in iter_block_items(doc):
        if isinstance(block, Paragraph) and block.text == "References":
            break
        paragraphs = [block] if isinstance(block, Paragraph) else [
            p for row in block.rows for cell in row.cells for p in cell.paragraphs
        ]
        for paragraph in paragraphs:
            for number in citation_numbers(paragraph.text):
                cited.add(number)
                if number not in ordered:
                    ordered.append(number)
    expected = list(range(1, len(refs) + 1))
    metrics["first_appearance_sequence"] = ordered
    if ordered != expected:
        errors.append(f"References are not first-appearance ordered: {ordered}")
    if cited != set(expected):
        errors.append(
            f"Citation/reference mismatch; missing={sorted(set(expected)-cited)}, "
            f"out_of_range={sorted(cited-set(expected))}"
        )

    doi_pattern = re.compile(r"https://doi\.org/([^\s.]+(?:\.[^\s.]+)*)", re.I)
    dois = []
    for paragraph in refs:
        match = doi_pattern.search(paragraph.text)
        if match:
            dois.append(match.group(1).rstrip("." ).lower())
    metrics["doi_count"] = len(dois)
    if len(dois) != len(set(dois)):
        errors.append("Duplicate DOI detected in the reference list")
    recent = sum(
        1
        for paragraph in refs
        if re.search(r"\((2023|2024|2025|2026)\)", paragraph.text)
        or re.search(r", (2023|2024|2025|2026)\.", paragraph.text)
    )
    metrics["references_2023_2026"] = recent
    if recent < 10:
        errors.append(f"Too few recent, directly relevant references: {recent}")

    forbidden = [
        "Fe(CN)6",
        "order-of-magnitude",
        "approximately 10×",
        "1.0×/1.5×/2.0×",
        "2.43×",
        "5.5/7",
        "3.8/7",
        "1.2/7",
    ]
    for phrase in forbidden:
        if phrase.lower() in text.lower():
            errors.append(f"Unsupported or obsolete phrase remains: {phrase}")

    required_phrases = [
        "Only Level 0 is currently supported",
        "no transfer coefficient",
        "not measured tBLG-SIB data",
        "Declaration of generative AI and AI-assisted technologies",
        "unchanged concept DOI",
        "https://doi.org/10.5281/zenodo.19362805",
        "The illustration is not a development sequence or readiness forecast",
        "Dashed links are non-directional interfaces",
        "Visual provenance map of the archived innovation-system materials",
    ]
    for phrase in required_phrases:
        if phrase.lower() not in text.lower():
            errors.append(f"Required evidence/disclosure phrase missing: {phrase}")

    if len(doc.tables) != 5:
        errors.append(f"Expected five editable manuscript tables, found {len(doc.tables)}")
    for index, table in enumerate(doc.tables):
        check_three_line(table, f"Table index {index}", errors)

    with ZipFile(MANUSCRIPT) as archive:
        document_xml = archive.read("word/document.xml")
        media = sorted(name for name in archive.namelist() if name.startswith("word/media/"))
        metrics["embedded_media"] = media
        if len(media) != 5:
            errors.append(f"Expected five embedded figures, found {len(media)}")
        # Match tracked-change elements themselves, not table-border elements
        # such as <w:insideH> and <w:insideV>.
        if re.search(rb"<w:(?:ins|del)(?=[\x20>])", document_xml):
            errors.append("Tracked insertions/deletions remain in the manuscript")
        if b"w:lnNumType" not in document_xml:
            errors.append("Continuous line-numbering settings are missing")
        if any("comments" in name.lower() for name in archive.namelist()):
            errors.append("Word comments remain in the manuscript")
        expected_figure_sizes = {
            1: (4500, 2640),
            2: (4000, 1675),
            3: (4000, 1540),
            4: (4500, 1886),
            5: (4500, 2180),
        }
        package_names = {
            1: "Figure_1_evidence_gates.png",
            2: "Figure_2_cost_trajectories_and_rank_correlations.png",
            3: "Figure_3_monte_carlo_cost_envelope.png",
            4: "Figure_4_integrated_assessment_chain.png",
            5: "Figure_5_tis_rubric_and_provenance.png",
        }
        for number in range(1, 6):
            name = f"word/media/image{number}.png"
            if name in archive.namelist():
                with archive.open(name) as stream:
                    image = Image.open(stream)
                    metrics[f"embedded_figure_{number}_pixels"] = list(image.size)
                    if image.size != expected_figure_sizes[number]:
                        errors.append(
                            f"Figure {number} has unexpected dimensions {image.size}; "
                            f"expected {expected_figure_sizes[number]}"
                        )
                package_figure = PACKAGE / "figures" / package_names[number]
                if not package_figure.exists():
                    errors.append(f"Regenerated package Figure {number} is missing")
                elif archive.read(name) != package_figure.read_bytes():
                    errors.append(f"Embedded Figure {number} differs from the regenerated package figure")

    highlights_doc = Document(HIGHLIGHTS)
    highlights = [p.text.strip() for p in highlights_doc.paragraphs if p.text.strip()]
    metrics["highlights"] = highlights
    if not 3 <= len(highlights) <= 5:
        errors.append(f"Highlights count outside 3–5: {len(highlights)}")
    for highlight in highlights:
        if len(highlight) > 85:
            errors.append(f"Highlight exceeds 85 characters: {highlight}")

    cover_text = all_text(Document(COVER))
    if EXPECTED_TITLE not in cover_text:
        errors.append("Cover letter contains a stale manuscript title")
    if "10.5281/zenodo.19362805" not in cover_text:
        errors.append("Cover letter does not preserve the unchanged Zenodo concept DOI")
    if "EST-D-26-05562" not in cover_text:
        errors.append("Cover letter does not identify the revision manuscript number")

    ga_text = pptx_text(GRAPHICAL_ABSTRACT)
    for obsolete in ("10×", "10x", "charge-transfer boost"):
        if obsolete.lower() in ga_text.lower():
            errors.append(f"Graphical abstract contains obsolete claim: {obsolete}")
    if "no validated battery transfer function" not in ga_text.lower():
        errors.append("Graphical abstract does not state the missing transfer function")

    response_text = all_text(Document(RESPONSE))
    for reviewer in ("Reviewer 1", "Reviewer 2", "Reviewer 3"):
        if reviewer not in response_text:
            errors.append(f"Point-by-point response is missing {reviewer}")

    mc = pd.read_csv(PACKAGE / "data" / "Monte_Carlo_TEA_2035_Cost.csv")
    total = mc["Total_System_Cost_USD_kWh"]
    expected_stats = {
        "median": round(float(total.median()), 1),
        "p5": round(float(total.quantile(0.05)), 1),
        "p95": round(float(total.quantile(0.95)), 1),
        "mean": round(float(total.mean()), 1),
    }
    metrics["recomputed_cost_stats"] = expected_stats
    numeric_text = text.replace("–", "-")
    for value in expected_stats.values():
        if f"{value:.1f}" not in numeric_text:
            errors.append(f"Recomputed cost statistic {value:.1f} is absent from manuscript")
    audit = pd.read_csv(PACKAGE / "data" / "TIS_Evidence_Audit.csv")
    if len(audit) != 7:
        errors.append(f"Evidence audit should have seven rows, found {len(audit)}")

    if not (PACKAGE / "figures" / "Figure_5_tis_rubric_and_provenance.png").exists():
        errors.append("Regenerated Figure 5 provenance audit is missing")
    if not (PACKAGE / "src" / "revise_submission.py").exists():
        errors.append("Reproducible manuscript-revision script is missing")
    if not (PACKAGE / "src" / "sync_submission_figures.py").exists():
        errors.append("Figure-synchronization script is missing")

    if "The corresponding author's phone number" not in (ROOT / "Submission_Readiness_Checklist.md").read_text(encoding="utf-8"):
        warnings.append("Checklist should remind the author to confirm a phone number")
    warnings.append("Publish the prepared Zenodo package as a new version before resubmission; the concept DOI remains unchanged.")
    warnings.append("Visual page-by-page inspection remains required after the final render.")

    result = {
        "status": "PASS" if not errors else "FAIL",
        "errors": sorted(set(errors)),
        "warnings": warnings,
        "metrics": metrics,
    }
    args.report.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
    print(json.dumps(result, indent=2, ensure_ascii=False))
    return 0 if not errors else 1


if __name__ == "__main__":
    raise SystemExit(main())
